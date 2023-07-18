import sys
import json

from PIL import Image
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.chart.data import XyChartData

typeToLayout = {'title': 0, 'text': 5, 'list': 1, 'picture': 5, 'plot': 5}


class Slide:
    def __init__(self, prs, slide_type, title, content, config=None):
        self.slide = prs.slides.add_slide(prs.slide_layouts[typeToLayout[slide_type]])
        self.prs = prs
        self.title = self.slide.shapes.title
        self.title.text = title
        self.content = content
        self.config = config

    def generate(self):
        pass


class TitleSlide(Slide):
    def generate(self):
        subtitle = self.slide.placeholders[1]
        subtitle.text = self.content


class TextSlide(Slide):
    def generate(self):
        left = width = height = Inches(1)
        top = Inches(2)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = self.content


class ListSlide(Slide):
    def generate(self):
        content_placeholder = self.slide.shapes.placeholders[1]
        for item in self.content:
            tf = content_placeholder.text_frame
            p = tf.add_paragraph()
            p.text = item['text']
            p.level = item['level']


class PictureSlide(Slide):
    def __init__(self, prs, slide_type, title, content, config):
        super().__init__(prs, slide_type, title, content, config)
        self.prs = prs

    def generate(self):
        img_path = self.content
        image = Image.open(img_path)
        width, height = image.size
        rescale_factor = width / height
        v_margin = Inches(2)
        h_margin = Inches(1)
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        target_width = max(slide_width - 2 * h_margin, h_margin)
        target_height = max(slide_height - 2 * v_margin, v_margin)
        if target_width / target_height <= rescale_factor:
            target_height = target_width / rescale_factor
            v_margin = 2 * (slide_height - target_height) / 3
        else:
            target_width = target_height * rescale_factor
            h_margin = (slide_width - target_width) / 2
        self.slide.shapes.add_picture(img_path, h_margin, v_margin, target_width, target_height)


class PlotSlide(Slide):
    def generate(self):
        chart_data = XyChartData()
        series_1 = chart_data.add_series('')
        with open(self.content) as file:
            for line in file:
                x, y = map(float, line.strip().split(';'))
                series_1.add_data_point(x, y)
        chart = self.slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS, Inches(1), Inches(2), Inches(6),
                                            Inches(4),
                                            chart_data).chart
        chart.has_legend = False
        chart.has_title = False
        y_axis = chart.value_axis
        y_axis.axis_title.text_frame.text = self.config['y-label']
        x_axis = chart.category_axis
        x_axis.axis_title.text_frame.text = self.config['x-label']


def generate_report(config_file):
    try:
        with open(config_file) as file:
            config = json.load(file)
    except json.JSONDecodeError as e:
        print(f'Invalid configuration file format: {config_file}')
        print(f'Error: {str(e)}')
        return

    prs = Presentation()

    slide_generators = {
        'title': TitleSlide,
        'text': TextSlide,
        'list': ListSlide,
        'picture': PictureSlide,
        'plot': PlotSlide,
    }

    for slide_data in config['presentation']:
        slide_type = slide_data['type']
        slide_generator = slide_generators.get(slide_type)
        if slide_generator:
            config = None
            if slide_type == 'plot':
                config = slide_data['configuration']
            slide = slide_generator(prs, slide_type, slide_data['title'], slide_data['content'], config)
            slide.generate()

    output_file_name = f'{config_file[:-5]}.pptx'
    prs.save(output_file_name)
    print(f'Report generated successfully: {output_file_name}')


if __name__ == '__main__':
    if len(sys.argv) != 2:
        print('Usage: python report_generator.py <config_file>')
    else:
        config_file = sys.argv[1]
        generate_report(config_file)
